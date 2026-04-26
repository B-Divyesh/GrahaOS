#!/usr/bin/env python3
"""
scripts/gen_ppt.py

Generates GrahaOS_Presentation.pptx at the project root — 19 slides for a
university presentation covering GrahaOS's novel ideas and general OS
concepts. Style: simple, academic, not flashy.

Run: python3 scripts/gen_ppt.py
Output: /home/atman/GrahaOS/GrahaOS_Presentation.pptx
"""
from __future__ import annotations

import pathlib

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------------------------------------------------------------------------
# Palette — muted, academic.
# ---------------------------------------------------------------------------
NAVY          = RGBColor(0x1B, 0x36, 0x5D)   # primary headings, rules
TEAL          = RGBColor(0x2E, 0x8B, 0x8B)   # accents, highlights
SLATE         = RGBColor(0x3C, 0x4A, 0x5B)   # body text
LIGHT_BG      = RGBColor(0xF5, 0xF7, 0xFA)   # page background
MEDIUM_GREY   = RGBColor(0x8A, 0x8F, 0x9B)   # secondary text
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG       = RGBColor(0xE8, 0xEE, 0xF4)
RULE_LIGHT    = RGBColor(0xCF, 0xD6, 0xDE)
SUCCESS_GREEN = RGBColor(0x2F, 0x7A, 0x4E)
WARN_AMBER    = RGBColor(0xC8, 0x82, 0x28)


# ---------------------------------------------------------------------------
# Canvas: 16:9 widescreen, 13.333 × 7.5 inches.
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def blank_slide():
    layout = prs.slide_layouts[6]  # fully blank
    slide = prs.slides.add_slide(layout)
    # Page background (off-white).
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def set_text(run, text, *, size=18, bold=False, color=SLATE,
             font_name="Calibri"):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text, *,
                 size=18, bold=False, color=SLATE, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_text(run, text, size=size, bold=bold, color=color,
             font_name=font_name)
    return tb


def add_bullet_list(slide, left, top, width, height, items, *,
                    size=18, color=SLATE, line_spacing=1.25,
                    bullet_color=TEAL):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # Colored bullet character, then body text.
        run_bullet = p.add_run()
        set_text(run_bullet, "•  ", size=size, bold=True,
                 color=bullet_color)
        run_body = p.add_run()
        set_text(run_body, item, size=size, color=color)
    return tb


def add_header(slide, title, subtitle=None):
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 prs.slide_width, Inches(0.10))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    # Title text
    add_text_box(slide, Inches(0.6), Inches(0.35), Inches(12.1),
                 Inches(0.75), title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(1.02), Inches(12.1),
                     Inches(0.45), subtitle, size=16, color=TEAL,
                     font_name="Calibri")

    # Thin rule below title
    rule = slide.shapes.add_connector(1, Inches(0.6), Inches(1.55),
                                      Inches(12.7), Inches(1.55))
    rule.line.color.rgb = RULE_LIGHT
    rule.line.width = Pt(1.25)


def add_footer(slide, page_num, total):
    # Page number and project name at bottom
    add_text_box(slide, Inches(0.6), Inches(7.10), Inches(4.0),
                 Inches(0.30),
                 "GrahaOS — Capability-First AI-Native OS",
                 size=10, color=MEDIUM_GREY)
    add_text_box(slide, Inches(11.0), Inches(7.10), Inches(1.8),
                 Inches(0.30), f"{page_num} / {total}",
                 size=10, color=MEDIUM_GREY, align=PP_ALIGN.RIGHT)


def box(slide, left, top, width, height, *, fill=CARD_BG, line=NAVY,
        line_width=Pt(1.0)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                 width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = line_width
    shp.shadow.inherit = False
    return shp


def rect(slide, left, top, width, height, *, fill=WHITE, line=NAVY,
         line_width=Pt(1.25)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width,
                                 height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = line_width
    shp.shadow.inherit = False
    return shp


def shape_text(shape, text, *, size=16, bold=False, color=SLATE,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf = shape.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    # Support multi-line via \n
    first = True
    for line in text.split("\n"):
        if first:
            run = p.add_run()
            first = False
        else:
            p = tf.add_paragraph()
            p.alignment = align
            run = p.add_run()
        set_text(run, line, size=size, bold=bold, color=color)


def arrow(slide, from_x, from_y, to_x, to_y, *, color=NAVY,
          width=Pt(1.5)):
    conn = slide.shapes.add_connector(2, from_x, from_y, to_x, to_y)
    conn.line.color.rgb = color
    conn.line.width = width
    # End arrow
    conn.line._get_or_add_ln()  # ensure <a:ln> exists
    return conn


TOTAL_SLIDES = 19


# ---------------------------------------------------------------------------
# Slide 1 — Title.
# ---------------------------------------------------------------------------
def slide_title():
    s = blank_slide()
    # Big navy block on left
    rect(s, 0, 0, Inches(4.5), prs.slide_height,
         fill=NAVY, line=NAVY)
    # Accent teal thin strip
    rect(s, Inches(4.5), 0, Inches(0.06), prs.slide_height,
         fill=TEAL, line=TEAL)

    # Title text on light side
    add_text_box(s, Inches(5.0), Inches(1.1), Inches(7.9),
                 Inches(1.4), "GrahaOS",
                 size=66, bold=True, color=NAVY)
    add_text_box(s, Inches(5.0), Inches(2.5), Inches(7.9),
                 Inches(0.55),
                 "A Capability-First, AI-Native Operating System",
                 size=22, color=TEAL)
    add_text_box(s, Inches(5.0), Inches(3.1), Inches(7.9),
                 Inches(0.5),
                 "Built from scratch  •  x86_64  •  Modular monolith → hybrid microkernel",
                 size=14, color=SLATE)

    # Horizontal rule
    rule = s.shapes.add_connector(1, Inches(5.0), Inches(3.85),
                                  Inches(12.8), Inches(3.85))
    rule.line.color.rgb = RULE_LIGHT
    rule.line.width = Pt(1)

    # Team block
    add_text_box(s, Inches(5.0), Inches(4.05), Inches(6.0),
                 Inches(0.4), "Team",
                 size=14, bold=True, color=TEAL)
    lines = [
        ("Divyesh Bine",   "100522733073",  "CSE Sem 8"),
        ("Farhan Ahmed",   "100522733079",  "CSE Sem 8"),
        ("Abraar Ahmed",   "100522733075",  "CSE Sem 8"),
    ]
    y = 4.45
    for name, roll, course in lines:
        add_text_box(s, Inches(5.0), Inches(y), Inches(3.6),
                     Inches(0.38), name, size=16, bold=True, color=NAVY)
        add_text_box(s, Inches(8.3), Inches(y), Inches(2.3),
                     Inches(0.38), roll, size=14, color=SLATE,
                     font_name="Consolas")
        add_text_box(s, Inches(10.8), Inches(y), Inches(2.0),
                     Inches(0.38), course, size=13, color=MEDIUM_GREY)
        y += 0.45

    # Vertical label on navy block
    add_text_box(s, Inches(0.6), Inches(3.3), Inches(3.5),
                 Inches(0.5), "University",
                 size=14, color=TEAL)
    add_text_box(s, Inches(0.6), Inches(3.75), Inches(3.5),
                 Inches(0.6), "Project Showcase",
                 size=26, bold=True, color=WHITE)
    add_text_box(s, Inches(0.6), Inches(5.3), Inches(3.5),
                 Inches(0.4),
                 "A research-grade OS written\nspec-first, test-driven",
                 size=12, color=RGBColor(0xBB, 0xCD, 0xE2))

    # Phase counter in corner
    add_text_box(s, Inches(0.6), Inches(6.7), Inches(3.5),
                 Inches(0.4),
                 "Phases 0–18 complete  •  399/399 tests",
                 size=10, color=RGBColor(0xBB, 0xCD, 0xE2))


# ---------------------------------------------------------------------------
# Slide 2 — Motivation.
# ---------------------------------------------------------------------------
def slide_motivation():
    s = blank_slide()
    add_header(s, "Why a New Operating System?",
               "The assumptions underneath Unix are 50 years old.")
    add_footer(s, 2, TOTAL_SLIDES)

    # Two-column layout
    add_text_box(s, Inches(0.6), Inches(1.85), Inches(6.0),
                 Inches(0.4),
                 "What current OSes were designed for",
                 size=18, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.6), Inches(2.3), Inches(6.0),
                    Inches(4.2), [
        "Multi-user timesharing on expensive hardware (1970s).",
        "Users, uid/gid, root — ambient authority.",
        "Fork + exec, signals, blocking read()/write().",
        "Files-as-bytes: no semantic meaning at OS level.",
        "Drivers welded into the kernel for speed.",
    ], size=16)

    add_text_box(s, Inches(7.0), Inches(1.85), Inches(5.7),
                 Inches(0.4),
                 "What modern workloads need",
                 size=18, bold=True, color=NAVY)
    add_bullet_list(s, Inches(7.0), Inches(2.3), Inches(5.7),
                    Inches(4.2), [
        "AI agents are a new class of user — must be sandboxable.",
        "Authority should be explicit, revocable, audit-traceable.",
        "Batched async I/O — syscalls cost too much per op.",
        "The OS should understand file semantics (type hashes, embeddings).",
        "Safe speculation: 'try this and roll back if wrong'.",
    ], size=16)

    # Tagline bar at bottom
    band = rect(s, Inches(0.6), Inches(6.55), Inches(12.1),
                Inches(0.5), fill=NAVY, line=NAVY)
    add_text_box(s, Inches(0.6), Inches(6.60), Inches(12.1),
                 Inches(0.42),
                 "GrahaOS rebuilds the kernel around capabilities, manifests, and speculation — native, not bolted-on.",
                 size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# Slide 3 — Seven Design Principles.
# ---------------------------------------------------------------------------
def slide_principles():
    s = blank_slide()
    add_header(s, "The Seven Design Principles",
               "Every spec file must uphold at least two of these and violate none.")
    add_footer(s, 3, TOTAL_SLIDES)

    principles = [
        ("1", "Capabilities are the sole authority",
         "No ambient auth. No users. Identity = tokens held."),
        ("2", "GCP is the only native public interface",
         "Every syscall manifest-typed. POSIX is optional userspace."),
        ("3", "Foundation before migration",
         "Lay new primitives, then move old callers behind shims."),
        ("4", "Failure is auditable, not silent",
         "Panics dump structured frames. Violations hit the audit log."),
        ("5", "Everything is testable by hand",
         "Every spec has manual_verification with exact commands."),
        ("6", "Drivers drift outward",
         "Kernel keeps VMM/sched/interrupts. Net/FS/block migrate out."),
        ("7", "Speculation is a first-class AI primitive",
         "Versioned FS + snapshots + txn channels = safe AI exploration."),
    ]
    col_w = Inches(5.85)
    row_h = Inches(1.15)
    x0 = Inches(0.6)
    x1 = Inches(6.9)
    for i, (num, title, detail) in enumerate(principles):
        col = i % 2
        row = i // 2
        x = x0 if col == 0 else x1
        y = Inches(1.80 + row * 1.25)

        b = box(s, x, y, col_w, row_h,
                fill=CARD_BG, line=NAVY, line_width=Pt(0.75))
        # Number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15),
                                    y + Inches(0.22), Inches(0.6),
                                    Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = NAVY
        circle.line.color.rgb = NAVY
        circle.shadow.inherit = False
        shape_text(circle, num, size=22, bold=True, color=WHITE)

        # Title + detail
        add_text_box(s, x + Inches(0.9), y + Inches(0.1), col_w - Inches(1.1),
                     Inches(0.5), title, size=15, bold=True, color=NAVY)
        add_text_box(s, x + Inches(0.9), y + Inches(0.55), col_w - Inches(1.1),
                     Inches(0.6), detail, size=12, color=SLATE)


# ---------------------------------------------------------------------------
# Slide 4 — Architecture.
# ---------------------------------------------------------------------------
def slide_architecture():
    s = blank_slide()
    add_header(s, "Architecture Overview",
               "Modular monolith today, hybrid microkernel by Phase 23.")
    add_footer(s, 4, TOTAL_SLIDES)

    # Userspace column
    us = rect(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.6),
              fill=WHITE, line=TEAL, line_width=Pt(1.5))
    add_text_box(s, Inches(0.8), Inches(1.65), Inches(5.5),
                 Inches(0.35), "USERSPACE", size=12, bold=True,
                 color=TEAL, align=PP_ALIGN.CENTER)

    us_boxes = [
        ("gash / grahai", "shell + AI agent"),
        ("ktest", "TAP test harness"),
        ("libc / libtap / libgcp", "thin libraries"),
        ("Future drivers (P21-23)", "net / AHCI / block"),
    ]
    for i, (title, sub) in enumerate(us_boxes):
        y = Inches(2.2 + i * 1.05)
        b = box(s, Inches(1.0), y, Inches(5.1), Inches(0.95),
                fill=CARD_BG, line=TEAL, line_width=Pt(0.5))
        add_text_box(s, Inches(1.15), y + Inches(0.1), Inches(4.8),
                     Inches(0.4), title, size=15, bold=True,
                     color=NAVY)
        add_text_box(s, Inches(1.15), y + Inches(0.5), Inches(4.8),
                     Inches(0.4), sub, size=12, color=SLATE)

    # Kernel column
    kn = rect(s, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.6),
              fill=WHITE, line=NAVY, line_width=Pt(1.5))
    add_text_box(s, Inches(7.0), Inches(1.65), Inches(5.5),
                 Inches(0.35), "KERNEL", size=12, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)

    kn_rows = [
        ("cap/ + io/ + ipc/ + mm/",
         "Tokens • Pledges • Channels • VMOs • Streams"),
        ("fs/ — VFS + GrahaFS + SimHash",
         "AI-metadata native, async_read hooks"),
        ("arch/x86_64/ — scheduler + VMM + PMM",
         "TASK_STATE_CHAN_WAIT, HHDM mapping, per-CPU"),
        ("drivers (in-kernel today, leaving P21-23)",
         "AHCI / e1000 / lapic / keyboard / fb"),
    ]
    for i, (title, sub) in enumerate(kn_rows):
        y = Inches(2.2 + i * 1.05)
        b = box(s, Inches(7.2), y, Inches(5.1), Inches(0.95),
                fill=CARD_BG, line=NAVY, line_width=Pt(0.5))
        add_text_box(s, Inches(7.35), y + Inches(0.1), Inches(4.8),
                     Inches(0.4), title, size=14, bold=True,
                     color=NAVY)
        add_text_box(s, Inches(7.35), y + Inches(0.5), Inches(4.8),
                     Inches(0.4), sub, size=12, color=SLATE)

    # Arrow between columns — drop into the gap between rows so it
    # doesn't land on top of a card's title.
    arrow(s, Inches(6.3), Inches(3.15), Inches(6.95), Inches(3.15),
          color=TEAL, width=Pt(2.25))
    add_text_box(s, Inches(6.18), Inches(2.80), Inches(0.95),
                 Inches(0.3), "IPC",
                 size=10, color=TEAL, align=PP_ALIGN.CENTER, bold=True)

    # Tagline
    add_text_box(s, Inches(0.6), Inches(6.75), Inches(12.1),
                 Inches(0.35),
                 "Kernel boundary moves outward each phase. Today it owns scheduler + MM + drivers. Phase 21-23: drivers out.",
                 size=13, color=SLATE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 5 — Capabilities.
# ---------------------------------------------------------------------------
def slide_capabilities():
    s = blank_slide()
    add_header(s, "Capabilities as the Sole Authority",
               "No users. No uid/gid. No root. Identity = bundle of tokens held.")
    add_footer(s, 5, TOTAL_SLIDES)

    # Left: token layout diagram
    add_text_box(s, Inches(0.6), Inches(1.85), Inches(6.3),
                 Inches(0.4), "cap_token_t — 64 bits",
                 size=18, bold=True, color=NAVY)

    # Visualize the token bit layout with three panels
    tx = Inches(0.6)
    ty = Inches(2.4)
    # generation 32 bits (wide)
    g = rect(s, tx, ty, Inches(3.0), Inches(1.2),
             fill=RGBColor(0xDC, 0xE8, 0xF4), line=NAVY, line_width=Pt(1.25))
    shape_text(g, "Generation\n32 bits", size=14, bold=True, color=NAVY)
    # idx 24 bits
    ix = rect(s, tx + Inches(3.0), ty, Inches(2.3), Inches(1.2),
              fill=RGBColor(0xCC, 0xE3, 0xDE), line=NAVY, line_width=Pt(1.25))
    shape_text(ix, "Object Index\n24 bits", size=14, bold=True, color=NAVY)
    # flags 8 bits
    fl = rect(s, tx + Inches(5.3), ty, Inches(1.0), Inches(1.2),
              fill=RGBColor(0xF0, 0xD6, 0xBE), line=NAVY, line_width=Pt(1.25))
    shape_text(fl, "Flags\n8", size=14, bold=True, color=NAVY)

    add_text_box(s, Inches(0.6), Inches(3.7), Inches(5.7),
                 Inches(0.4),
                 "Passed by value in registers. Never heap-allocated.",
                 size=12, color=MEDIUM_GREY)

    # Properties
    add_text_box(s, Inches(0.6), Inches(4.2), Inches(6.3),
                 Inches(0.4), "Resolve path (hot)",
                 size=15, bold=True, color=TEAL)
    add_bullet_list(s, Inches(0.6), Inches(4.6), Inches(6.3),
                    Inches(2.4), [
        "Atomic-acquire load generation → compare → reject if stale.",
        "Check calling pid in audience_set (CAP_FLAG_PUBLIC bypass).",
        "Bitmask AND check on rights (READ / WRITE / DERIVE / …).",
        "No locks. ~20–30 cycles typical.",
    ], size=13)

    # Right: rights + revocation
    add_text_box(s, Inches(7.3), Inches(1.85), Inches(5.5),
                 Inches(0.4), "Key properties",
                 size=18, bold=True, color=NAVY)

    props = [
        ("Derivable",
         "Narrow rights / audience into child caps; tree tracked."),
        ("Revocable",
         "Single gen++ invalidates all outstanding tokens (O(1))."),
        ("Audience-scoped",
         "Audience set is per-capability; privacy-preserving filter."),
        ("Kind-dispatched",
         "CHANNEL, VMO, STREAM, FILE, PROC — each has a deactivator."),
    ]
    for i, (lbl, txt) in enumerate(props):
        y = Inches(2.4 + i * 1.1)
        b = box(s, Inches(7.3), y, Inches(5.4), Inches(0.95),
                fill=CARD_BG, line=TEAL, line_width=Pt(0.5))
        add_text_box(s, Inches(7.45), y + Inches(0.1), Inches(5.1),
                     Inches(0.4), lbl, size=14, bold=True,
                     color=NAVY)
        add_text_box(s, Inches(7.45), y + Inches(0.48), Inches(5.1),
                     Inches(0.42), txt, size=12, color=SLATE)


# ---------------------------------------------------------------------------
# Slide 6 — Pledges.
# ---------------------------------------------------------------------------
def slide_pledges():
    s = blank_slide()
    add_header(s, "Pledges — The Second Authority Axis",
               "'What WILL this process ever do again?' — monotonically narrowing.")
    add_footer(s, 6, TOTAL_SLIDES)

    # Left: concept
    add_text_box(s, Inches(0.6), Inches(1.85), Inches(6.0),
                 Inches(0.4), "Concept",
                 size=18, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.6), Inches(2.3), Inches(6.0),
                    Inches(3.2), [
        "Capability answers: CAN do this now?",
        "Pledge answers: WILL ever do this again?",
        "Narrow-only: drop PLEDGE_FS_WRITE once → permanent.",
        "Reset only on SYS_SPAWN (new child, fresh subset).",
        "Every sensitive syscall = one AND + one compare.",
    ], size=14)

    # Right: 12 class grid
    add_text_box(s, Inches(7.0), Inches(1.85), Inches(5.7),
                 Inches(0.4), "12 Classes (16-bit mask)",
                 size=18, bold=True, color=NAVY)

    classes = [
        ("FS_READ", "files read"),
        ("FS_WRITE", "files write"),
        ("NET_CLIENT", "outbound net"),
        ("NET_SERVER", "listen/bind"),
        ("SPAWN", "create processes"),
        ("IPC_SEND", "channel send"),
        ("IPC_RECV", "channel recv"),
        ("SYS_QUERY", "read kernel state"),
        ("SYS_CONTROL", "reboot / kill"),
        ("AI_CALL", "AI agent"),
        ("COMPUTE", "CPU + memory"),
        ("TIME", "sleep / clock"),
    ]
    x0 = Inches(7.0)
    y0 = Inches(2.35)
    cw = Inches(1.85)
    ch = Inches(0.7)
    for i, (name, hint) in enumerate(classes):
        col = i % 3
        row = i // 3
        x = x0 + col * cw
        y = y0 + row * Inches(0.8)
        b = box(s, x + Inches(0.05), y, cw - Inches(0.1), ch,
                fill=CARD_BG, line=TEAL, line_width=Pt(0.5))
        add_text_box(s, x + Inches(0.1), y + Inches(0.05), cw - Inches(0.2),
                     Inches(0.3), name, size=12, bold=True, color=NAVY)
        add_text_box(s, x + Inches(0.1), y + Inches(0.34), cw - Inches(0.2),
                     Inches(0.3), hint, size=10, color=SLATE)

    # Bottom note
    band = rect(s, Inches(0.6), Inches(6.45), Inches(12.1),
                Inches(0.55), fill=CARD_BG, line=NAVY, line_width=Pt(0.75))
    add_text_box(s, Inches(0.85), Inches(6.55), Inches(11.6),
                 Inches(0.38),
                 "pledge_allows(task, CLASS)  =  (mask.raw & (1u << CLASS)) != 0   — header-only inline, zero overhead.",
                 size=12, color=SLATE, font_name="Consolas")


# ---------------------------------------------------------------------------
# Slide 7 — CAN.
# ---------------------------------------------------------------------------
def slide_can():
    s = blank_slide()
    add_header(s, "CAN — Capability Activation Network",
               "The sole registry of system components. Drivers as graph nodes.")
    add_footer(s, 7, TOTAL_SLIDES)

    add_bullet_list(s, Inches(0.6), Inches(1.85), Inches(5.8),
                    Inches(4.5), [
        "Every driver / service / app registered as a CAN cap.",
        "Dependency edges: e.g., tcp_ip depends on e1000 depends on pci_bus.",
        "Activate cascades forward; deactivate cascades backward.",
        "Each cap has on_activate / on_deactivate callbacks — physically drives hardware (keyboard IRQ on/off, fb mmio on/off).",
        "Audit log entries on every state flip.",
        "Phase 16 completed the hardware wiring (keyboard, fb, e1000, ahci).",
    ], size=14)

    # Right: simplified graph
    # Diagram: cpu/memory at bottom, pci, e1000, tcp_ip, then apps on top
    gx = Inches(7.2)
    # Layer boxes
    layers = [
        (Inches(6.15), "apps + services", [("gash", 0.2), ("grahai", 1.4), ("ktest", 2.6)]),
        (Inches(5.05), "FS + net + audit", [("grahafs", 0.2), ("tcp_ip", 1.4), ("audit", 2.6)]),
        (Inches(3.95), "drivers", [("keyboard", 0.1), ("fb", 1.1), ("e1000", 1.9), ("ahci", 2.9)]),
        (Inches(2.85), "hardware", [("cpu", 0.2), ("memory", 1.2), ("pci_bus", 2.3)]),
    ]
    for y_inch, label, boxes_ in layers:
        add_text_box(s, gx, y_inch - Inches(0.05), Inches(1.4),
                     Inches(0.25), label, size=10, color=TEAL, bold=True)
        for name, off_x in boxes_:
            b = box(s, gx + Inches(1.5 + off_x), y_inch,
                    Inches(1.0), Inches(0.55),
                    fill=WHITE, line=NAVY, line_width=Pt(0.75))
            shape_text(b, name, size=10, color=NAVY, bold=True)

    # Up arrows
    for sx_off, ex_off in [(1.7, 1.7), (2.9, 2.9), (4.0, 4.0)]:
        arrow(s, gx + Inches(sx_off), Inches(3.3),
              gx + Inches(ex_off), Inches(3.05),
              color=MEDIUM_GREY, width=Pt(1.0))
        arrow(s, gx + Inches(sx_off), Inches(4.45),
              gx + Inches(ex_off), Inches(4.15),
              color=MEDIUM_GREY, width=Pt(1.0))
        arrow(s, gx + Inches(sx_off), Inches(5.55),
              gx + Inches(ex_off), Inches(5.25),
              color=MEDIUM_GREY, width=Pt(1.0))

    add_text_box(s, gx + Inches(1.4), Inches(6.6), Inches(4.7),
                 Inches(0.35),
                 "→ activate cascades up; deactivate cascades down.",
                 size=11, color=SLATE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 8 — GCP.
# ---------------------------------------------------------------------------
def slide_gcp():
    s = blank_slide()
    add_header(s, "GCP — The Only Native Public Interface",
               "/etc/gcp.json is the single source of truth for every syscall and message.")
    add_footer(s, 8, TOTAL_SLIDES)

    # Left: concept
    add_text_box(s, Inches(0.6), Inches(1.85), Inches(5.8),
                 Inches(0.4), "What GCP enforces",
                 size=18, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.6), Inches(2.3), Inches(5.8),
                    Inches(3.2), [
        "Every syscall has a manifest entry in JSON.",
        "Every channel message carries a 64-bit type_hash (FNV-1a).",
        "Kernel rejects mistyped sends at entry (−EPROTOTYPE).",
        "Every stream op has a generated op_schema row.",
        "AI agents + humans read the same JSON file.",
        "POSIX is optional libposix in userspace — never in kernel.",
    ], size=14)

    # Right: JSON sample
    code_bg = rect(s, Inches(7.0), Inches(1.85), Inches(5.8),
                   Inches(4.2), fill=RGBColor(0x1F, 0x29, 0x37),
                   line=NAVY, line_width=Pt(0.75))
    add_text_box(s, Inches(7.0), Inches(1.85), Inches(5.8),
                 Inches(0.32), "  /etc/gcp.json  (excerpt)",
                 size=11, color=RGBColor(0x8A, 0xC9, 0xE2),
                 font_name="Consolas")
    sample = ("\"grahaos.io.v1\": {\n"
              "  \"ops\": {\n"
              "    \"READ_VMO\": { \"op\": 1,\n"
              "      \"pledge\": \"FS_READ\",\n"
              "      \"dispatcher\": \"dispatch_read_vmo\" },\n"
              "    \"WRITE_VMO\":  { \"op\": 2, … },\n"
              "    \"SENDMSG\":    { \"op\": 3, … }\n"
              "  }\n"
              "}")
    code_tb = s.shapes.add_textbox(Inches(7.2), Inches(2.25),
                                   Inches(5.5), Inches(3.6))
    ctf = code_tb.text_frame
    ctf.word_wrap = True
    first = True
    for line in sample.split("\n"):
        p = ctf.paragraphs[0] if first else ctf.add_paragraph()
        first = False
        run = p.add_run()
        set_text(run, line, size=13, color=RGBColor(0xE2, 0xE8, 0xEF),
                 font_name="Consolas")

    # Bottom callout
    band = rect(s, Inches(0.6), Inches(6.3), Inches(12.1),
                Inches(0.7), fill=NAVY, line=NAVY)
    add_text_box(s, Inches(0.75), Inches(6.35), Inches(11.8),
                 Inches(0.3),
                 "gen_manifest.py → single JSON → kernel op_schema + userspace OP_* constants.",
                 size=13, bold=True, color=WHITE)
    add_text_box(s, Inches(0.75), Inches(6.68), Inches(11.8),
                 Inches(0.3),
                 "Kernel and userspace can never drift out of sync — the ABI is generated.",
                 size=12, color=RGBColor(0xBB, 0xCD, 0xE2))


# ---------------------------------------------------------------------------
# Slide 9 — Channels.
# ---------------------------------------------------------------------------
def slide_channels():
    s = blank_slide()
    add_header(s, "Channels — Typed Bidirectional IPC",
               "Replace ambient pipes. Atomic handle transfer. Manifest-anchored.")
    add_footer(s, 9, TOTAL_SLIDES)

    # Diagram: two processes, channel between them
    # Process A
    pa = box(s, Inches(0.7), Inches(2.2), Inches(2.8), Inches(2.4),
             fill=CARD_BG, line=NAVY, line_width=Pt(1.25))
    add_text_box(s, Inches(0.7), Inches(2.3), Inches(2.8),
                 Inches(0.35), "Process A", size=14, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
    hA = box(s, Inches(1.0), Inches(3.1), Inches(2.2), Inches(0.55),
             fill=WHITE, line=TEAL, line_width=Pt(0.6))
    shape_text(hA, "wr_handle", size=11, color=NAVY)
    add_text_box(s, Inches(0.7), Inches(3.85), Inches(2.8),
                 Inches(0.6),
                 "Holds CAP_KIND_CHANNEL\nwrite endpoint.",
                 size=10, color=SLATE, align=PP_ALIGN.CENTER)

    # Process B
    pb = box(s, Inches(9.9), Inches(2.2), Inches(2.8), Inches(2.4),
             fill=CARD_BG, line=NAVY, line_width=Pt(1.25))
    add_text_box(s, Inches(9.9), Inches(2.3), Inches(2.8),
                 Inches(0.35), "Process B", size=14, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
    hB = box(s, Inches(10.2), Inches(3.1), Inches(2.2), Inches(0.55),
             fill=WHITE, line=TEAL, line_width=Pt(0.6))
    shape_text(hB, "rd_handle", size=11, color=NAVY)
    add_text_box(s, Inches(9.9), Inches(3.85), Inches(2.8),
                 Inches(0.6),
                 "Holds CAP_KIND_CHANNEL\nread endpoint.",
                 size=10, color=SLATE, align=PP_ALIGN.CENTER)

    # Channel ring in middle
    ch = box(s, Inches(4.0), Inches(2.2), Inches(5.6), Inches(2.4),
             fill=RGBColor(0xE6, 0xEE, 0xE1), line=NAVY, line_width=Pt(1.5))
    add_text_box(s, Inches(4.0), Inches(2.3), Inches(5.6),
                 Inches(0.35), "channel_t (manifest-typed)",
                 size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(4.0), Inches(2.65), Inches(5.6),
                 Inches(0.3), "type_hash = FNV-1a(\"grahaos.io.v1\")",
                 size=10, color=TEAL, align=PP_ALIGN.CENTER,
                 font_name="Consolas")
    # slot grid
    for i in range(8):
        sx = Inches(4.15 + i * 0.68)
        slot = rect(s, sx, Inches(3.2), Inches(0.6), Inches(0.75),
                    fill=WHITE, line=NAVY, line_width=Pt(0.4))
        shape_text(slot, f"s{i}", size=9, color=SLATE)
    add_text_box(s, Inches(4.0), Inches(4.05), Inches(5.6),
                 Inches(0.5),
                 "Fixed 320 B slots  •  hdr + 256 B inline + 8 × handle",
                 size=11, color=SLATE, align=PP_ALIGN.CENTER)

    # Arrows
    arrow(s, Inches(3.55), Inches(3.4), Inches(3.95), Inches(3.4),
          color=TEAL, width=Pt(2.0))
    arrow(s, Inches(9.65), Inches(3.4), Inches(10.15), Inches(3.4),
          color=TEAL, width=Pt(2.0))

    # Bottom bullets
    add_bullet_list(s, Inches(0.7), Inches(5.0), Inches(12.0),
                    Inches(2.0), [
        "Two CAP_KIND_CHANNEL cap_objects per channel — read endpoint + write endpoint.",
        "chan_send with mismatched type_hash → −EPROTOTYPE + audit entry.",
        "Atomic handle transfer: handles moved out of sender's table into the slot, then into receiver.",
        "TASK_STATE_CHAN_WAIT + per-tick deadline scan give timeouts without calibration.",
    ], size=13)


# ---------------------------------------------------------------------------
# Slide 10 — VMOs.
# ---------------------------------------------------------------------------
def slide_vmos():
    s = blank_slide()
    add_header(s, "Virtual Memory Objects (VMOs)",
               "Named, refcounted physical frames. Copy-on-write — the speculation substrate.")
    add_footer(s, 10, TOTAL_SLIDES)

    # Diagram: parent VMO → clone → pp_refcount on shared pages
    # Parent
    pv = box(s, Inches(0.7), Inches(2.0), Inches(3.4), Inches(3.2),
             fill=CARD_BG, line=NAVY, line_width=Pt(1.5))
    add_text_box(s, Inches(0.7), Inches(2.1), Inches(3.4),
                 Inches(0.35), "Parent VMO", size=14, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
    for i, label in enumerate(["page 0", "page 1", "page 2", "page 3"]):
        sx = Inches(0.9 + (i % 2) * 1.55)
        sy = Inches(2.6 + (i // 2) * 1.1)
        p = rect(s, sx, sy, Inches(1.35), Inches(0.9),
                 fill=WHITE, line=NAVY, line_width=Pt(0.5))
        shape_text(p, label, size=11, color=NAVY)

    # Child (COW clone)
    cv = box(s, Inches(9.2), Inches(2.0), Inches(3.4), Inches(3.2),
             fill=RGBColor(0xFA, 0xEF, 0xDD), line=NAVY, line_width=Pt(1.5))
    add_text_box(s, Inches(9.2), Inches(2.1), Inches(3.4),
                 Inches(0.35), "COW Child VMO", size=14, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
    for i, label in enumerate(["page 0 (RO)", "page 1 (RO)", "page 2 (RO)", "page 3 (RO)"]):
        sx = Inches(9.4 + (i % 2) * 1.55)
        sy = Inches(2.6 + (i // 2) * 1.1)
        p = rect(s, sx, sy, Inches(1.35), Inches(0.9),
                 fill=WHITE, line=NAVY, line_width=Pt(0.5))
        shape_text(p, label, size=10, color=NAVY)

    # Shared frame in middle
    shared = rect(s, Inches(5.5), Inches(3.0), Inches(2.3), Inches(1.3),
                  fill=RGBColor(0xE6, 0xEE, 0xE1), line=TEAL,
                  line_width=Pt(1.5))
    shape_text(shared, "Shared physical\nframe\npp_refcount = 2",
               size=12, color=NAVY, bold=True)
    arrow(s, Inches(4.2), Inches(3.55), Inches(5.45), Inches(3.55),
          color=TEAL, width=Pt(1.75))
    arrow(s, Inches(9.15), Inches(3.55), Inches(7.85), Inches(3.55),
          color=TEAL, width=Pt(1.75))

    # Bottom bullets
    add_text_box(s, Inches(0.7), Inches(5.4), Inches(12.0),
                 Inches(0.35), "Why VMOs matter",
                 size=16, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.7), Inches(5.75), Inches(12.0),
                    Inches(1.5), [
        "Zero-copy shared memory between kernel and userspace (HHDM direct-map).",
        "COW isolation: write to a shared frame triggers private-copy via page-fault handler.",
        "Per-physical-page refcount (1 byte parallel to PMM bitmap).",
        "Substrate for Phase 25 speculative transactions — \"try and roll back\" for free.",
    ], size=13)


# ---------------------------------------------------------------------------
# Slide 11 — Submission Streams.
# ---------------------------------------------------------------------------
def slide_streams():
    s = blank_slide()
    add_header(s, "Submission Streams — Async I/O",
               "io_uring-style batching, GCP-native. Every op manifest-validated.")
    add_footer(s, 11, TOTAL_SLIDES)

    # Diagram: userspace → SQ VMO → kernel worker → CQ VMO → userspace
    us = box(s, Inches(0.6), Inches(2.5), Inches(2.3), Inches(2.0),
             fill=CARD_BG, line=NAVY, line_width=Pt(1.25))
    shape_text(us, "Userspace\nproducer",
               size=13, bold=True, color=NAVY)

    sq = box(s, Inches(3.3), Inches(2.0), Inches(2.6), Inches(1.5),
             fill=RGBColor(0xDC, 0xE8, 0xF4), line=NAVY,
             line_width=Pt(1.25))
    shape_text(sq, "SQ VMO\n(sqe_t 64 B × N)",
               size=12, bold=True, color=NAVY)

    w = box(s, Inches(6.3), Inches(2.5), Inches(2.4), Inches(2.0),
            fill=RGBColor(0xFA, 0xEF, 0xDD), line=NAVY,
            line_width=Pt(1.25))
    shape_text(w, "Kernel\nworker\nthread",
               size=13, bold=True, color=NAVY)

    cq = box(s, Inches(9.1), Inches(2.0), Inches(2.6), Inches(1.5),
             fill=RGBColor(0xE6, 0xEE, 0xE1), line=NAVY,
             line_width=Pt(1.25))
    shape_text(cq, "CQ VMO\n(cqe_t 32 B × 2N)",
               size=12, bold=True, color=NAVY)

    us2 = box(s, Inches(12.0), Inches(2.5), Inches(0.8), Inches(2.0),
              fill=CARD_BG, line=NAVY, line_width=Pt(0.75))
    shape_text(us2, "User\nreaps", size=10, bold=True, color=NAVY)

    # Arrows
    arrow(s, Inches(2.95), Inches(2.75), Inches(3.3), Inches(2.75),
          color=TEAL, width=Pt(2.0))
    add_text_box(s, Inches(2.8), Inches(2.25), Inches(0.7),
                 Inches(0.3), "SQEs", size=10, color=TEAL,
                 bold=True, align=PP_ALIGN.CENTER)

    arrow(s, Inches(5.95), Inches(2.75), Inches(6.3), Inches(2.75),
          color=NAVY, width=Pt(1.75))
    arrow(s, Inches(8.75), Inches(2.75), Inches(9.1), Inches(2.75),
          color=NAVY, width=Pt(1.75))
    add_text_box(s, Inches(8.35), Inches(2.25), Inches(0.9),
                 Inches(0.3), "CQEs", size=10, color=NAVY, bold=True,
                 align=PP_ALIGN.CENTER)

    arrow(s, Inches(11.75), Inches(2.75), Inches(12.05), Inches(2.75),
          color=TEAL, width=Pt(2.0))

    # Memory-layout caption — move below all ring boxes so it doesn't
    # overlap the kernel worker tile.
    add_text_box(s, Inches(0.6), Inches(4.6), Inches(12.2),
                 Inches(0.3),
                 "head @ offset 0   |   tail @ offset 128   |   entries from page 1   (cache-line separated)",
                 size=11, color=MEDIUM_GREY, align=PP_ALIGN.CENTER,
                 font_name="Consolas")

    # Bullets — shift down to make room for the caption above.
    add_bullet_list(s, Inches(0.6), Inches(5.05), Inches(12.1),
                    Inches(2.0), [
        "One SYS_STREAM_SUBMIT can hand off thousands of ops — amortizes user/kernel crossings.",
        "Every SQE type-validated against the stream's manifest at kernel entry.",
        "Per-op pledge check; bad ops get inline CQE with −EPROTOTYPE or −EPLEDGE.",
        "ACQUIRE / RELEASE on head/tail — io_uring memory discipline, documented in stream.h.",
        "Optional notify channel: non-blocking event loop without polling.",
    ], size=13)


# ---------------------------------------------------------------------------
# Slide 12 — GrahaFS + SimHash.
# ---------------------------------------------------------------------------
def slide_grahafs():
    s = blank_slide()
    add_header(s, "GrahaFS — Filesystem with AI Semantics",
               "AI metadata and similarity clustering are first-class.")
    add_footer(s, 12, TOTAL_SLIDES)

    # Left: on-disk layout
    add_text_box(s, Inches(0.6), Inches(1.85), Inches(5.8),
                 Inches(0.4), "On-disk structure",
                 size=18, bold=True, color=NAVY)
    rows = [
        ("Superblock",      "magic, sizes, free bitmap"),
        ("Inode table",     "256-byte inodes (not 128)"),
        ("Extended meta",   "tags, importance, ai_embedding[128]"),
        ("Data blocks",     "4 KiB each; 48 KiB file cap → Phase 19 fixes"),
    ]
    for i, (lbl, hint) in enumerate(rows):
        y = Inches(2.3 + i * 0.9)
        b = box(s, Inches(0.6), y, Inches(5.8), Inches(0.8),
                fill=CARD_BG, line=NAVY, line_width=Pt(0.5))
        add_text_box(s, Inches(0.8), y + Inches(0.1), Inches(5.4),
                     Inches(0.35), lbl, size=14, bold=True, color=NAVY)
        add_text_box(s, Inches(0.8), y + Inches(0.43), Inches(5.4),
                     Inches(0.35), hint, size=12, color=SLATE)

    # Right: SimHash clustering
    add_text_box(s, Inches(6.8), Inches(1.85), Inches(6.0),
                 Inches(0.4), "SimHash Clustering (Phase 11)",
                 size=18, bold=True, color=NAVY)
    add_bullet_list(s, Inches(6.8), Inches(2.3), Inches(6.0),
                    Inches(3.4), [
        "64-bit fingerprint per file (FNV-1a over 3-char shingles).",
        "Sequential-leader clustering: Hamming distance ≤ 10 → same cluster.",
        "Cluster IDs stored in inode's ai_reserved[].",
        "Background indexer task scans root every 3 s.",
        "/illu/ directory for deterministic files (no auto-cluster).",
    ], size=13)

    # Bottom: visual of bit distance
    band = rect(s, Inches(0.6), Inches(6.25), Inches(12.1),
                Inches(0.75), fill=NAVY, line=NAVY)
    add_text_box(s, Inches(0.85), Inches(6.30), Inches(11.6),
                 Inches(0.32),
                 "simhash(\"hello world this is a test\") = 0x7F2A4B…  |  popcount(a ^ b) = Hamming distance",
                 size=12, color=WHITE, font_name="Consolas")
    add_text_box(s, Inches(0.85), Inches(6.63), Inches(11.6),
                 Inches(0.32),
                 "OS-level similarity — no external ML dependency.",
                 size=11, color=RGBColor(0xBB, 0xCD, 0xE2))


# ---------------------------------------------------------------------------
# Slide 13 — AI first-class.
# ---------------------------------------------------------------------------
def slide_ai():
    s = blank_slide()
    add_header(s, "AI as a First-Class Citizen",
               "Not bolted on. Baked in at the filesystem and speculation layers.")
    add_footer(s, 13, TOTAL_SLIDES)

    cards = [
        ("grahai",
         "Embedded AI agent process. Calls Gemini via in-kernel TCP/IP stack (Phase 9e).",
         "PLEDGE_AI_CALL + audit log entry per invoke."),
        ("ai_metadata on every file",
         "tags[], importance score, ai_embedding[128], last_ai_touched.",
         "Queryable via SYS_SET_AI_METADATA / SYS_GET_AI_METADATA."),
        ("SimHash similarity",
         "Find near-duplicate files in Hamming distance — no external embedding service.",
         "OS surfaces cluster membership as metadata."),
        ("Speculative execution  (Phase 25)",
         "Transactions on COW VMO snapshots. AI tries a change; OS rolls back cleanly if wrong.",
         "This is the differentiator — every other OS nails 'safety' but not safe exploration."),
    ]
    for i, (title, desc, tag) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.25)
        y = Inches(2.0 + row * 2.45)
        b = box(s, x, y, Inches(6.0), Inches(2.25),
                fill=CARD_BG, line=NAVY, line_width=Pt(0.75))
        add_text_box(s, x + Inches(0.25), y + Inches(0.15),
                     Inches(5.5), Inches(0.45),
                     title, size=16, bold=True, color=NAVY)
        add_text_box(s, x + Inches(0.25), y + Inches(0.6),
                     Inches(5.5), Inches(1.0),
                     desc, size=13, color=SLATE)
        add_text_box(s, x + Inches(0.25), y + Inches(1.65),
                     Inches(5.5), Inches(0.5),
                     tag, size=11, color=TEAL)


# ---------------------------------------------------------------------------
# Slide 14 — Spec-Driven Workflow.
# ---------------------------------------------------------------------------
def slide_workflow():
    s = blank_slide()
    add_header(s, "The Spec-Driven Workflow",
               "Every phase starts as a YAML spec. No speculative coding.")
    add_footer(s, 14, TOTAL_SLIDES)

    # Process chain
    steps = [
        ("1", "YAML spec",
         "specs/phase-NN.yml\nactor workflows\nmanual verification"),
        ("2", "Implement",
         "new source files\nstruct + syscall\nmanifests"),
        ("3", "TAP tests",
         "make test\ngate_tests_passing\nstability iters"),
        ("4", "Manual VP",
         "shell commands\nexact expected\noutputs match"),
        ("5", "Audit",
         "problems_faced.md\naudit_log.json\nper-issue root cause"),
        ("6", "Gate next",
         "only start N+1 after\nN exit-criteria green"),
    ]
    for i, (n, title, body) in enumerate(steps):
        x = Inches(0.6 + i * 2.1)
        y = Inches(2.2)
        # Step card
        b = box(s, x, y, Inches(1.9), Inches(2.6),
                fill=CARD_BG, line=NAVY, line_width=Pt(0.75))
        # Number circle
        c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7),
                               y + Inches(0.15), Inches(0.55),
                               Inches(0.55))
        c.fill.solid()
        c.fill.fore_color.rgb = NAVY
        c.line.color.rgb = NAVY
        c.shadow.inherit = False
        shape_text(c, n, size=18, bold=True, color=WHITE)
        add_text_box(s, x + Inches(0.05), y + Inches(0.8),
                     Inches(1.8), Inches(0.35), title,
                     size=13, bold=True, color=NAVY,
                     align=PP_ALIGN.CENTER)
        add_text_box(s, x + Inches(0.1), y + Inches(1.15),
                     Inches(1.7), Inches(1.4), body,
                     size=10, color=SLATE, align=PP_ALIGN.CENTER)
        # Arrow to next
        if i < len(steps) - 1:
            arrow(s, x + Inches(1.92), y + Inches(1.3),
                  x + Inches(2.08), y + Inches(1.3),
                  color=TEAL, width=Pt(1.5))

    # Callout
    band = rect(s, Inches(0.6), Inches(5.6), Inches(12.1),
                Inches(1.3), fill=NAVY, line=NAVY)
    add_text_box(s, Inches(0.9), Inches(5.75), Inches(11.6),
                 Inches(0.4),
                 "Why this works",
                 size=14, bold=True, color=TEAL)
    add_text_box(s, Inches(0.9), Inches(6.15), Inches(11.6),
                 Inches(0.75),
                 "Every phase has a reviewable YAML contract before a line of C is written. "
                 "Every audit entry in problems_faced.md becomes an invariant for future phases — "
                 "we know why the code is the way it is.",
                 size=13, color=WHITE)


# ---------------------------------------------------------------------------
# Slide 15 — Testing discipline.
# ---------------------------------------------------------------------------
def slide_testing():
    s = blank_slide()
    add_header(s, "Testing & Verification Floor",
               "TAP 1.4. Every spec has a manual_verification section. make test exits zero or fails.")
    add_footer(s, 15, TOTAL_SLIDES)

    # Metrics strip
    metrics = [
        ("399/399",  "assertions passing",  SUCCESS_GREEN),
        ("31",       "TAP test files",      NAVY),
        ("15/15",    "stability iterations clean", SUCCESS_GREEN),
        ("26s",      "full suite wall time", TEAL),
        ("9",        "issues logged in phase 18", WARN_AMBER),
    ]
    for i, (big, small, color) in enumerate(metrics):
        x = Inches(0.6 + i * 2.5)
        y = Inches(2.0)
        b = box(s, x, y, Inches(2.35), Inches(1.4),
                fill=WHITE, line=color, line_width=Pt(1.25))
        add_text_box(s, x, y + Inches(0.15), Inches(2.35),
                     Inches(0.6), big,
                     size=28, bold=True, color=color,
                     align=PP_ALIGN.CENTER)
        add_text_box(s, x, y + Inches(0.85), Inches(2.35),
                     Inches(0.4), small,
                     size=11, color=SLATE, align=PP_ALIGN.CENTER)

    # How tests are structured
    add_text_box(s, Inches(0.6), Inches(3.75), Inches(6.0),
                 Inches(0.4), "Structure per phase",
                 size=15, bold=True, color=NAVY)
    add_bullet_list(s, Inches(0.6), Inches(4.2), Inches(6.0),
                    Inches(2.7), [
        "unit_tests — isolated property checks.",
        "integration_tests — full-system workflows.",
        "fault_injection_tests — deliberate failures.",
        "stress_tests — soak at scale.",
        "gate: yes  vs  gate: no per entry.",
    ], size=13)

    add_text_box(s, Inches(7.1), Inches(3.75), Inches(5.6),
                 Inches(0.4), "Notable tests",
                 size=15, bold=True, color=NAVY)
    add_bullet_list(s, Inches(7.1), Inches(4.2), Inches(5.6),
                    Inches(2.7), [
        "sentinel_meta proves the harness catches failures.",
        "cantest_v2 (37) / captest_v2 (42) — capability lifecycle.",
        "chantest (24) / vmotest (18) — Phase 17 IPC substrate.",
        "streamtest (22) — Phase 18 async I/O correctness.",
        "canstress 100× cap cycle — no stuck IRQs.",
    ], size=13)


# ---------------------------------------------------------------------------
# Slide 16 — Phases completed.
# ---------------------------------------------------------------------------
def slide_completed():
    s = blank_slide()
    add_header(s, "Phases Completed (0 → 18)",
               "Boot to Phase 18 async I/O — 18 substantial spec-driven milestones.")
    add_footer(s, 16, TOTAL_SLIDES)

    # Left: early phases list
    early = [
        ("0–7",  "Bootloader, GDT/IDT, PMM/VMM, sched"),
        ("8",    "CAN registry + drivers (kb/fb/lapic)"),
        ("9",    "E1000 + Mongoose TCP/IP + HTTPS + DNS + Gemini"),
        ("10",   "Per-process FDs + pipes + redirects"),
        ("11",   "SimHash + sequential-leader clustering"),
    ]
    add_text_box(s, Inches(0.6), Inches(1.85), Inches(5.8),
                 Inches(0.4), "Foundation",
                 size=16, bold=True, color=TEAL)
    for i, (p, d) in enumerate(early):
        y = Inches(2.25 + i * 0.85)
        b = box(s, Inches(0.6), y, Inches(5.8), Inches(0.75),
                fill=CARD_BG, line=NAVY, line_width=Pt(0.5))
        add_text_box(s, Inches(0.8), y + Inches(0.05),
                     Inches(0.8), Inches(0.35),
                     f"P{p}", size=13, bold=True, color=NAVY)
        add_text_box(s, Inches(1.7), y + Inches(0.16),
                     Inches(4.0), Inches(0.45),
                     d, size=12, color=SLATE)

    # Right: recent phases
    late = [
        ("12", "Test harness + TAP 1.4 floor"),
        ("13", "Structured logging + panic/oops dump"),
        ("14", "Slab + kheap + per-CPU (NUMA-ready)"),
        ("15a", "Capability Objects v2 (tokens, gen, audience)"),
        ("15b", "Pledges + persistent audit log"),
        ("16", "CAN callbacks → physical hardware control"),
        ("17", "Channels + VMOs + handle-passing IPC"),
        ("18", "Submission Streams (async I/O)  ← current"),
    ]
    add_text_box(s, Inches(7.0), Inches(1.85), Inches(5.8),
                 Inches(0.4), "Recent — spec-driven cycle",
                 size=16, bold=True, color=TEAL)
    for i, (p, d) in enumerate(late):
        y = Inches(2.25 + i * 0.55)
        b = box(s, Inches(7.0), y, Inches(5.8), Inches(0.48),
                fill=CARD_BG if i < len(late) - 1 else RGBColor(0xD4, 0xE5, 0xD8),
                line=NAVY, line_width=Pt(0.5))
        add_text_box(s, Inches(7.2), y + Inches(0.06),
                     Inches(0.9), Inches(0.3),
                     f"P{p}", size=13, bold=True, color=NAVY)
        add_text_box(s, Inches(8.2), y + Inches(0.08),
                     Inches(4.4), Inches(0.32),
                     d, size=12,
                     color=SUCCESS_GREEN if i == len(late) - 1 else SLATE,
                     bold=(i == len(late) - 1))


# ---------------------------------------------------------------------------
# Slide 17 — Roadmap 19-28.
# ---------------------------------------------------------------------------
def slide_roadmap():
    s = blank_slide()
    add_header(s, "Roadmap 19 → 28",
               "Toward an MVP: versioned FS, userspace drivers, speculation, Wasmtime, TUI.")
    add_footer(s, 17, TOTAL_SLIDES)

    phases = [
        ("19", "GrahaFS v2",
         "versioned segments + full data journaling; async_read consumed"),
        ("20", "Per-CPU scheduler",
         "work-stealing runqueues; resource limits"),
        ("21", "Userspace driver framework",
         "e1000 migration out of kernel"),
        ("22", "In-tree TCP/IP",
         "replace 19 k-line Mongoose; ARP/IPv4/UDP/TCP-Reno/DHCP"),
        ("23", "AHCI userspace",
         "interrupt-driven block I/O; completes drivers-outward"),
        ("24", "COW snapshot subsystem",
         "whole-process snapshots on VMO infra"),
        ("25", "Speculative execution",
         "transactions on snapshots — the AI-exploration primitive"),
        ("26", "Wasmtime userspace",
         "GCP as WIT; Rust user programs"),
        ("27", "TUI framework",
         "virtual consoles + widgets"),
        ("28", "MVP",
         "TUI-first shell + AI polish + 1-hour stress bench"),
    ]
    col_w = Inches(6.1)
    row_h = Inches(0.55)
    for i, (n, title, desc) in enumerate(phases):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.2)
        y = Inches(2.0 + row * 0.87)
        b = box(s, x, y, col_w, row_h,
                fill=CARD_BG, line=NAVY, line_width=Pt(0.5))
        # Phase number badge
        badge = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x + Inches(0.05), y, Inches(0.7), row_h)
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.color.rgb = TEAL
        badge.shadow.inherit = False
        shape_text(badge, f"P{n}", size=14, bold=True, color=WHITE)
        add_text_box(s, x + Inches(0.85), y + Inches(0.04),
                     Inches(5.0), Inches(0.28),
                     title, size=13, bold=True, color=NAVY)
        add_text_box(s, x + Inches(0.85), y + Inches(0.28),
                     Inches(5.0), Inches(0.28),
                     desc, size=10, color=SLATE)

    # Bottom callout — Phase 25 highlight
    band = rect(s, Inches(0.6), Inches(6.8), Inches(12.1),
                Inches(0.55), fill=TEAL, line=TEAL)
    add_text_box(s, Inches(0.8), Inches(6.88), Inches(11.7),
                 Inches(0.4),
                 "Phase 25 is the thesis: snapshots + COW VMO + transactional channels = safe AI speculation.",
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 18 — Key Novel Contributions.
# ---------------------------------------------------------------------------
def slide_contributions():
    s = blank_slide()
    add_header(s, "Key Novel Contributions",
               "What's actually new — vs. Unix / Linux / microkernels / io_uring.")
    add_footer(s, 18, TOTAL_SLIDES)

    items = [
        ("No users, no root",
         "Identity = bundle of cap_tokens. uid/gid/root never existed."),
        ("Dual authority (cap + pledge)",
         "Caps say CAN now. Pledges say WILL ever again. Orthogonal."),
        ("GCP as sole public interface",
         "Every syscall and message validates against JSON manifest at entry."),
        ("Manifest-typed channels",
         "Ambient-authority pipes dead. Every message type-hash-anchored."),
        ("Speculation as OS primitive",
         "COW VMOs + snapshots + transactional channels (coming P24–25)."),
        ("Generated ABI — kernel ↔ user single-source",
         "gen_manifest.py emits both sides. Drift impossible."),
        ("AI metadata in the filesystem",
         "tags, importance, embeddings, SimHash clusters — native."),
        ("Everything auditable",
         "Caps, pledges, stream rejections, VMO faults — all in /var/audit."),
    ]
    for i, (title, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.2)
        y = Inches(1.9 + row * 1.25)
        b = box(s, x, y, Inches(6.0), Inches(1.1),
                fill=CARD_BG, line=NAVY, line_width=Pt(0.75))
        # Accent diamond on the left
        diamond = s.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                     x + Inches(0.2),
                                     y + Inches(0.32),
                                     Inches(0.45), Inches(0.45))
        diamond.fill.solid()
        diamond.fill.fore_color.rgb = TEAL
        diamond.line.color.rgb = TEAL
        diamond.shadow.inherit = False
        add_text_box(s, x + Inches(0.8), y + Inches(0.12),
                     Inches(5.0), Inches(0.45),
                     title, size=14, bold=True, color=NAVY)
        add_text_box(s, x + Inches(0.8), y + Inches(0.55),
                     Inches(5.0), Inches(0.55),
                     body, size=11, color=SLATE)


# ---------------------------------------------------------------------------
# Slide 19 — Conclusion.
# ---------------------------------------------------------------------------
def slide_conclusion():
    s = blank_slide()
    add_header(s, "Conclusion",
               "An AI-native OS built the slow way — spec first, test floor high.")
    add_footer(s, 19, TOTAL_SLIDES)

    # Big takeaways
    items = [
        ("Where we are",
         "Phases 0–18 complete. 399/399 tests. Async I/O shipped. Kernel ready to shed drivers."),
        ("What's next",
         "Phase 19 versioned FS, then drivers-outward (P21–23), then speculation (P24–25)."),
        ("Why it matters",
         "Every other OS treats AI as an app. We treat AI as a user class with its own safety needs."),
    ]
    for i, (title, body) in enumerate(items):
        y = Inches(1.95 + i * 1.2)
        b = box(s, Inches(0.7), y, Inches(11.9), Inches(1.05),
                fill=CARD_BG, line=NAVY, line_width=Pt(0.75))
        add_text_box(s, Inches(0.95), y + Inches(0.1), Inches(11.5),
                     Inches(0.4), title,
                     size=16, bold=True, color=NAVY)
        add_text_box(s, Inches(0.95), y + Inches(0.5), Inches(11.5),
                     Inches(0.55), body,
                     size=13, color=SLATE)

    # Q&A band
    qa = rect(s, Inches(0.7), Inches(5.8), Inches(11.9), Inches(1.2),
              fill=NAVY, line=NAVY)
    add_text_box(s, Inches(0.7), Inches(5.95), Inches(11.9),
                 Inches(0.5), "Thank you",
                 size=30, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.7), Inches(6.55), Inches(11.9),
                 Inches(0.4),
                 "Questions?  —  github.com/B-Divyesh/GrahaOS",
                 size=14, color=RGBColor(0xBB, 0xCD, 0xE2),
                 align=PP_ALIGN.CENTER)


# Register slides in order.
slide_title()
slide_motivation()
slide_principles()
slide_architecture()
slide_capabilities()
slide_pledges()
slide_can()
slide_gcp()
slide_channels()
slide_vmos()
slide_streams()
slide_grahafs()
slide_ai()
slide_workflow()
slide_testing()
slide_completed()
slide_roadmap()
slide_contributions()
slide_conclusion()


OUT = pathlib.Path(__file__).resolve().parent.parent / "GrahaOS_Presentation.pptx"
prs.save(str(OUT))
print(f"wrote {OUT}  —  {len(prs.slides)} slides")
